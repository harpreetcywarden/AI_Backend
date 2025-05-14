import os
import csv  # Added for CSV support
from langchain.text_splitter import RecursiveCharacterTextSplitter
from sentence_transformers import SentenceTransformer
from langchain_community.document_loaders import PyPDFLoader
from docx import Document as DocxDocument  # For .docx files
from pptx import Presentation             # For .pptx files
from PIL import Image                     # For image processing
import pytesseract                        # For OCR (text extraction from images)
from pdf2image import convert_from_path   # For extracting images from PDFs
import pandas as pd                       # For Excel file processing

# Set the path to the Tesseract executable
pytesseract.pytesseract.tesseract_cmd = "/usr/bin/tesseract"

def process_pdf(file_path):
    """Extract text and images from a PDF file."""
    try:
        # Step 1: Extract text from PDF pages
        loader = PyPDFLoader(file_path)
        documents = loader.load()
        pdf_text = [doc.page_content.strip() for doc in documents if doc.page_content.strip()]

        # Step 2: Extract images from PDF pages
        images = convert_from_path(file_path)  # Convert PDF pages to images
        ocr_text = []
        for image in images:
            text = pytesseract.image_to_string(image).strip()  # Perform OCR on each image
            if text:
                ocr_text.append(text)

        # Step 3: Combine text from PDF pages and OCR results
        combined_text = pdf_text + ocr_text
        return combined_text if combined_text else []  # Return combined text
    except Exception as e:
        print(f"⚠️ Error processing PDF {file_path}: {str(e)}")
        return []

def process_txt(file_path):
    """Extract text from a plain text file."""
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read().strip()
    return [text] if text else []  # Ensure empty files don't return empty chunks

def process_image(file_path):
    """Extract text from an image using OCR."""
    try:
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image).strip()
        print(f"Extracted Text from Image: {text}")  # Debugging
        return [text] if text else []  # Avoid storing empty OCR results
    except pytesseract.pytesseract.TesseractNotFoundError:
        print("⚠️ Tesseract is not installed or not in PATH. Please install Tesseract.")
        return []
    except Exception as e:
        print(f"⚠️ Error processing image {file_path}: {str(e)}")
        return []

def process_docx(file_path):
    """Extract text from a Word document."""
    doc = DocxDocument(file_path)
    text = "\n".join([para.text.strip() for para in doc.paragraphs if para.text.strip()])
    return [text] if text else []  # Ensure only non-empty paragraphs are stored

def process_pptx(file_path):
    """Extract text from a PowerPoint presentation."""
    presentation = Presentation(file_path)
    slides_text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slides_text.append(shape.text.strip())

    return slides_text if slides_text else []  # Avoid storing empty slides

def process_xlsx(file_path):
    """Extract data from an Excel file."""
    try:
        # Read the Excel file
        excel_data = pd.ExcelFile(file_path)
        sheets = {sheet_name: excel_data.parse(sheet_name) for sheet_name in excel_data.sheet_names}

        # Convert each sheet's data to text
        text_chunks = []
        for sheet_name, df in sheets.items():
            # Convert the DataFrame to a string representation
            text_chunks.append(f"Sheet: {sheet_name}\n{df.to_string(index=False)}")

        return text_chunks if text_chunks else []  # Return text chunks
    except Exception as e:
        print(f"⚠️ Error processing Excel file {file_path}: {str(e)}")
        return []

def process_csv(file_path):
    """Extract data from a CSV file."""
    try:
        text_chunks = []
        with open(file_path, 'r', encoding='utf-8') as csvfile:
            # Read CSV and convert each row to a string
            csv_reader = csv.reader(csvfile)
            header = next(csv_reader)  # Get header row
            for i, row in enumerate(csv_reader):
                # Combine header and row values into a readable string
                row_text = ", ".join([f"{col}: {val}" for col, val in zip(header, row)])
                text_chunks.append(f"Row {i+1}: {row_text}")
        
        return text_chunks if text_chunks else []
    except Exception as e:
        print(f"⚠️ Error processing CSV file {file_path}: {str(e)}")
        return []

